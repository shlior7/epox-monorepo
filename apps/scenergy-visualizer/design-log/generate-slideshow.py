#!/usr/bin/env python3
"""
Generate a PPTX slideshow for Scenergy Visualizer UI flows.
Run: python generate-slideshow.py
Output: Scenergy-Visualizer-UI-Slideshow.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BG = RGBColor(24, 24, 27)  # zinc-900
PRIMARY = RGBColor(99, 102, 241)  # indigo-500
ACCENT = RGBColor(34, 197, 94)  # green-500
TEXT_WHITE = RGBColor(255, 255, 255)
TEXT_GRAY = RGBColor(161, 161, 170)  # zinc-400
CARD_BG = RGBColor(39, 39, 42)  # zinc-800


def add_title_slide(title, subtitle=""):
    """Add a title slide with dark theme."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(54)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = TEXT_WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        tf.paragraphs[0].text = subtitle
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.color.rgb = TEXT_GRAY
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(title, content_items, has_columns=False):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = CARD_BG
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = TEXT_WHITE
    
    # Content
    if has_columns and len(content_items) == 2:
        # Two columns
        for i, (col_title, items) in enumerate(content_items):
            left = Inches(0.5) if i == 0 else Inches(6.8)
            
            # Column title
            col_box = slide.shapes.add_textbox(left, Inches(1.6), Inches(5.5), Inches(0.5))
            tf = col_box.text_frame
            tf.paragraphs[0].text = col_title
            tf.paragraphs[0].font.size = Pt(22)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = PRIMARY
            
            # Items
            content_box = slide.shapes.add_textbox(left, Inches(2.2), Inches(5.5), Inches(4.5))
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for j, item in enumerate(items):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(18)
                p.font.color.rgb = TEXT_WHITE
                p.space_after = Pt(12)
    else:
        # Single column
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            if isinstance(item, tuple):
                # Header + content
                p.text = item[0]
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = PRIMARY
                p.space_after = Pt(4)
                
                for sub in item[1]:
                    p = tf.add_paragraph()
                    p.text = f"    • {sub}"
                    p.font.size = Pt(18)
                    p.font.color.rgb = TEXT_WHITE
                    p.space_after = Pt(8)
            else:
                p.text = f"• {item}"
                p.font.size = Pt(20)
                p.font.color.rgb = TEXT_WHITE
                p.space_after = Pt(12)
    
    return slide


def add_diagram_slide(title, diagram_text):
    """Add a slide with ASCII diagram."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = CARD_BG
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = TEXT_WHITE
    
    # Diagram box
    diagram_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5)
    )
    diagram_shape.fill.solid()
    diagram_shape.fill.fore_color.rgb = CARD_BG
    diagram_shape.line.fill.background()
    
    # Diagram text
    diagram_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(5))
    tf = diagram_box.text_frame
    tf.word_wrap = False
    tf.paragraphs[0].text = diagram_text
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.name = "Courier New"
    tf.paragraphs[0].font.color.rgb = TEXT_WHITE
    
    return slide


def add_table_slide(title, headers, rows):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = CARD_BG
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = TEXT_WHITE
    
    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1
    col_width = Inches(12 / num_cols)
    
    table = slide.shapes.add_table(
        num_rows, num_cols, 
        Inches(0.5), Inches(1.6), 
        Inches(12), Inches(0.5 * num_rows)
    ).table
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_WHITE
            p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_flow_cards_slide(title, cards):
    """Add a slide with flow cards."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = CARD_BG
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = TEXT_WHITE
    
    # Cards
    card_width = Inches(2.8)
    card_height = Inches(4.5)
    start_left = Inches(0.5)
    gap = Inches(0.3)
    
    for i, (card_title, card_items) in enumerate(cards):
        left = start_left + (card_width + gap) * i
        
        # Card shape
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = PRIMARY
        card.line.width = Pt(2)
        
        # Card title
        title_shape = slide.shapes.add_textbox(left + Inches(0.1), Inches(1.8), card_width - Inches(0.2), Inches(0.5))
        tf = title_shape.text_frame
        tf.paragraphs[0].text = card_title
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = PRIMARY
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Card content
        content_shape = slide.shapes.add_textbox(left + Inches(0.15), Inches(2.4), card_width - Inches(0.3), Inches(3.5))
        tf = content_shape.text_frame
        tf.word_wrap = True
        
        for j, item in enumerate(card_items):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_WHITE
            p.space_after = Pt(6)
        
        # Arrow between cards (except last)
        if i < len(cards) - 1:
            arrow_left = left + card_width + Inches(0.05)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arrow_left, Inches(3.5), Inches(0.2), Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = PRIMARY
            arrow.line.fill.background()
    
    return slide


# ============================================================
# GENERATE SLIDES
# ============================================================

# Slide 1: Title
add_title_slide(
    "🎨 Scenergy Visualizer",
    "AI-Powered Product Visualization Platform"
)

# Slide 2: Overview
add_content_slide(
    "What is Scenergy Visualizer?",
    [
        "AI-powered platform for generating professional product images",
        "Transform product photos into styled room visualizations",
        "Bulk generation for entire catalogs (50-1000+ products)",
        "Direct sync to e-commerce stores (Shopify, WooCommerce, BigCommerce)",
        "WebP-optimized images ready for CDN delivery"
    ]
)

# Slide 3: The Problem
add_content_slide(
    "The Problem We're Solving",
    [
        ("Traditional Product Photography", [
            "Expensive: $500-1,000+ per product",
            "Slow: 3-5 weeks turnaround",
            "Limited scene variety",
            "Manual coordination required"
        ]),
        ("Our Solution", [
            "Generate hundreds of images in minutes",
            "Fraction of the cost",
            "Unlimited style variations",
            "Fully automated workflow"
        ])
    ]
)

# Slide 4: Two User Paths
add_content_slide(
    "Two Ways to Get Started",
    [
        ("🏪 Path A: Connect Store", [
            "Connect Shopify, WooCommerce, or BigCommerce",
            "Import products directly",
            "Approve images → Auto-sync to store",
            "Full bidirectional sync with ERP ID"
        ]),
        ("📤 Path B: Upload Manually", [
            "No store connection required",
            "Upload product images directly",
            "Generate images",
            "Download only (no store sync)"
        ])
    ],
    has_columns=True
)

# Slide 5: User Flow
add_flow_cards_slide(
    "User Journey: Store Connected",
    [
        ("1. Connect", ["OAuth with store", "Select platform", "Authorize access"]),
        ("2. Import", ["By product IDs", "By category", "Or all products"]),
        ("3. Generate", ["Select products", "Choose inspiration", "AI generates images"]),
        ("4. Review", ["Side-by-side compare", "Approve / Reject", "Edit if needed"]),
    ]
)

# Slide 6: Product Import
add_table_slide(
    "Product Import Options",
    ["Method", "Best For", "Description"],
    [
        ["By Product IDs", "Specific items", "Paste SKUs or product IDs"],
        ["By Category", "Focused batches", "Select store categories"],
        ["All Products", "Full catalog", "Import everything at once"],
    ]
)

# Slide 7: Subscription Tiers
add_table_slide(
    "Subscription Tiers",
    ["Plan", "Max Products", "Monthly Credits", "Price"],
    [
        ["Starter", "50", "100", "Free"],
        ["Pro", "100", "1,000", "$49/mo"],
        ["Business", "500", "5,000", "$199/mo"],
        ["Enterprise", "1,000+", "Unlimited", "Custom"],
    ]
)

# Slide 8: Generation Workflow
add_content_slide(
    "Collection Session Workflow",
    [
        "Step 1: Select Products — Choose products from your library",
        "Step 2: AI Analysis — AI detects product types, styles, colors",
        "Step 3: Choose Inspiration — Upload or search for room images",
        "Step 4: Generate — AI creates styled visualizations",
        "Step 5: Review & Approve — Side-by-side comparison",
        "Step 6: Sync to Store — Approved images auto-sync"
    ]
)

# Slide 9: Review Experience
add_content_slide(
    "Review & Approval Experience",
    [
        ("Side-by-Side Comparison", [
            "Original product image vs Generated scene",
            "Drag slider to compare",
            "Zoom sync between images"
        ]),
        ("Review Actions", [
            "✓ Approve — Mark for store sync (free)",
            "✎ Edit — Open editor modal (free)",
            "↻ Regenerate — New version (1 credit)",
            "✗ Reject — Keep in library, no sync (free)"
        ]),
        ("Bulk Review Mode", [
            "Keyboard shortcuts: A (approve), R (reject), E (edit)",
            "Arrow keys to navigate",
            "Progress indicator: '15 of 45 reviewed'"
        ])
    ]
)

# Slide 10: Credit Actions
add_table_slide(
    "Credit System",
    ["Action", "Credits", "Notes"],
    [
        ["Generate image", "1", "Per image variation"],
        ["Regenerate", "1", "New version with same/new settings"],
        ["Approve", "0", "Free — marks for sync"],
        ["Edit", "0", "Free — opens editor"],
        ["Download", "0", "Free — always available"],
        ["Reject", "0", "Free — stays in library"],
    ]
)

# Slide 11: Imported vs Uploaded
add_table_slide(
    "Imported vs Uploaded Products",
    ["Feature", "Imported", "Uploaded"],
    [
        ["Generate images", "✓", "✓"],
        ["Edit images", "✓", "✓"],
        ["Download images", "✓", "✓"],
        ["Favorite / Tag", "✓", "✓"],
        ["Approve for Store", "✓", "✗"],
        ["Sync to Store", "✓", "✗"],
        ["ERP ID tracking", "✓", "✗"],
    ]
)

# Slide 12: Key Screens
add_content_slide(
    "Key Screens",
    [
        ("Dashboard", ["Recent generations", "Quick stats", "Pending review items"]),
        ("Products Library", ["Searchable catalog", "Source badges (Imported/Uploaded)", "Quick actions"]),
        ("Collection Wizard", ["4-step guided flow", "Product selection → Settings → Generation"]),
        ("Review Gallery", ["Grid view with status filters", "Bulk actions", "Approval workflow"]),
        ("Store Sync Dashboard", ["Connection status", "Sync queue", "History"]),
    ]
)

# Slide 13: Product Analysis
add_content_slide(
    "AI Product Analysis",
    [
        "Products are analyzed to understand:",
        "• Product type (sofa, desk, lamp, etc.)",
        "• Materials (leather, wood, metal, fabric)",
        "• Color palette (primary and accent colors)",
        "• Style (modern, rustic, minimalist)",
        "• Best room types (living room, bedroom, office)",
        "",
        "This analysis improves prompt engineering for better results!"
    ]
)

# Slide 14: Store Sync
add_content_slide(
    "Store Sync Workflow",
    [
        "1. Connect Store — One-time OAuth setup (Shopify, WooCommerce, BigCommerce)",
        "2. Import Products — Products get ERP ID for tracking",
        "3. Generate Images — Normal generation workflow",
        "4. Review & Approve — Mark images as 'Approved for Store'",
        "5. Auto-Sync — Approved images automatically sync to product pages",
        "6. Track Status — View sync history, pending, failed items"
    ]
)

# Slide 15: Key Features
add_content_slide(
    "✨ Key Features",
    [
        "🔗 Bidirectional Sync — Import from store, export images back",
        "🤖 AI-Powered Analysis — Better prompts from product understanding",
        "👁️ Side-by-Side Review — Easy comparison for quality control",
        "⚡ Bulk Generation — 50+ products at once with consistent style",
        "📦 WebP Optimization — Images ready for CDN delivery",
        "⌨️ Keyboard Shortcuts — Fast review with A/R/E keys",
        "💳 Transparent Credits — Always know the cost before action"
    ]
)

# Slide 16: Mobile Support
add_content_slide(
    "📱 Mobile Responsive",
    [
        "Dashboard — Stats cards, quick actions",
        "Products — Grid view, swipe navigation",
        "Review — Full-screen images, tap actions",
        "Bottom navigation — Home, Products, Generate, Settings"
    ]
)

# Slide 17: Next Steps
add_content_slide(
    "Roadmap",
    [
        ("Now: MVP (Q1 2026)", [
            "Collection Sessions (bulk generation)",
            "Standalone Generation Flows",
            "WebP optimization",
            "Approval workflow",
            "Store sync preparation"
        ]),
        ("Next: Q2 2026", [
            "Full store sync (Shopify, WooCommerce, BigCommerce)",
            "In-image editing (mask and regenerate)",
            "Video generation (product turntables)"
        ]),
        ("Later: Q3-Q4 2026", [
            "API access for developers",
            "3D model generation",
            "Advanced analytics dashboard"
        ])
    ]
)

# Slide 18: End
add_title_slide(
    "🎨 Scenergy Visualizer",
    "Questions? Contact the Product Team\n\nSee Design Logs #001-#008 for technical details"
)

# Save the presentation
output_path = "Scenergy-Visualizer-UI-Slideshow.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")

